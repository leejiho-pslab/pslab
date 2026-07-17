#!/usr/bin/env python3
# AWK 에어워크주니어 — 편집 가능한 PowerPoint (밝은 테마)
import os
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

AWK="/home/user/pslab/assets/awk"
CON="/home/user/pslab/assets/awk-content"
SP="/tmp/claude-0/-home-user-pslab/e949e4d9-ac40-56b3-bb03-0a51ce53f94b/scratchpad"

WHITE=RGBColor(0xFF,0xFF,0xFF)
SOFT =RGBColor(0xF5,0xF7,0xFB)
INK  =RGBColor(0x14,0x14,0x14)
INK2 =RGBColor(0x5A,0x5A,0x66)
INK3 =RGBColor(0x9A,0x9A,0xA6)
BLUE =RGBColor(0x1F,0x5F,0xFF)
CORAL=RGBColor(0xFF,0x5A,0x1F)
LINE =RGBColor(0xE4,0xE8,0xF0)
KFONT="Pretendard"
SW,SH=13.333,7.5
PAD=0.7

prs=Presentation(); prs.slide_width=I(SW); prs.slide_height=I(SH)
BLANK=prs.slide_layouts[6]

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    return s

def setp(p,txt,size,color,bold=False,font=KFONT,align=PP_ALIGN.LEFT,spacing=1.0):
    p.alignment=align
    if spacing: p.line_spacing=spacing
    run=p.add_run(); run.text=txt
    f=run.font; f.size=Pt(size); f.bold=bold; f.name=font; f.color.rgb=color
    return run

def text(s,l,t,w,h,txt,size,color,bold=False,font=KFONT,align=PP_ALIGN.LEFT,spacing=1.06,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(I(l),I(t),I(w),I(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0; tf.vertical_anchor=anchor
    first=True
    for ln in txt.split("\n"):
        p=tf.paragraphs[0] if first else tf.add_paragraph()
        setp(p,ln,size,color,bold,font,align,spacing); first=False
    return tb

def kicker(s,txt,color=BLUE,l=PAD,t=0.6):
    ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,I(l),I(t+0.08),I(0.28),I(0.05))
    ln.fill.solid(); ln.fill.fore_color.rgb=CORAL; ln.line.fill.background(); ln.shadow.inherit=False
    text(s,l+0.4,t,10,0.35,txt,12.5,color,bold=True)

def brand(s):
    text(s,PAD,0.32,6,0.3,"AWK · 에어워크주니어",11.5,INK,bold=True)

def rrect(s,l,t,w,h,fill=WHITE,border=LINE,bw=0.75):
    c=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,I(l),I(t),I(w),I(h))
    c.fill.solid(); c.fill.fore_color.rgb=fill
    if border is None: c.line.fill.background()
    else: c.line.color.rgb=border; c.line.width=Pt(bw)
    c.shadow.inherit=False
    a=c.adjustments
    try: a[0]=0.06
    except: pass
    return c

def card(s,l,t,w,h,title,body,n=None,ncolor=BLUE,tsize=16,bsize=12.5,fill=WHITE):
    c=rrect(s,l,t,w,h,fill=fill)
    tf=c.text_frame; tf.word_wrap=True
    tf.margin_left=I(0.2);tf.margin_right=I(0.2);tf.margin_top=I(0.18);tf.margin_bottom=I(0.16)
    first=True
    if n:
        p=tf.paragraphs[0]; setp(p,n,11,ncolor,True,KFONT,PP_ALIGN.LEFT,1.0); first=False
    if title:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); setp(p,title,tsize,INK,True,KFONT,PP_ALIGN.LEFT,1.1); first=False
    if body:
        p=tf.add_paragraph(); p.space_before=Pt(4); setp(p,body,bsize,INK2,False,KFONT,PP_ALIGN.LEFT,1.25)
    return c

def pic_cover(s,path,l,t,w,h):
    """지정 박스를 채우는 cover-crop 이미지."""
    im=Image.open(path); iw,ih=im.size
    boxar=h/w; imar=ih/iw
    if imar>boxar:  # 이미지가 더 세로 → 위아래 크롭
        cw=iw; ch=int(iw*boxar); y=(ih-ch)//2; crop=(0,y,cw,y+ch)
    else:
        ch=ih; cw=int(ih/boxar); x=(iw-cw)//2; crop=(x,0,x+cw,ih)
    im2=im.crop(crop)
    tmp=f"{SP}/_pptmp_{os.path.basename(path)}_{int(w*100)}x{int(h*100)}.jpg"
    im2.convert("RGB").save(tmp,"JPEG",quality=90)
    s.shapes.add_picture(tmp,I(l),I(t),I(w),I(h))

def pic_fit(s,path,l,t,w,h):
    """박스 안에 비율 유지로 넣고, 실제 크기 반환(가운데 정렬)."""
    im=Image.open(path); iw,ih=im.size; ar=ih/iw
    W=w; H=w*ar
    if H>h: H=h; W=h/ar
    x=l+(w-W)/2; y=t+(h-H)/2
    s.shapes.add_picture(path,I(x),I(y),I(W),I(H))
    return x,y,W,H

def pill(s,l,t,txt,fill,fg=WHITE,size=11):
    w=0.16+len(txt)*0.09
    c=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,I(l),I(t),I(w),I(0.34))
    c.fill.solid(); c.fill.fore_color.rgb=fill; c.line.fill.background(); c.shadow.inherit=False
    c.adjustments[0]=0.5
    tf=c.text_frame; tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    setp(tf.paragraphs[0],txt,size,fg,True,KFONT,PP_ALIGN.CENTER,1.0)
    return w

# ============ SLIDES ============

# 1 COVER
s=slide()
pic_cover(s, f"{AWK}/awk-011.jpg", SW-5.6,0,5.6,SH)
# left white panel already white bg
text(s,PAD,0.34,7,0.3,"CONTENT REFERENCE · PSLAB × AWK",11,BLUE,bold=True)
text(s,PAD,2.5,7.2,2.2,"AWK 에어워크주니어\n브랜드 & 콘텐츠 제안",40,INK,True,KFONT,spacing=1.08)
# 코랄 언더바
u=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,I(PAD),I(4.35),I(1.0),I(0.09)); u.fill.solid(); u.fill.fore_color.rgb=CORAL; u.line.fill.background(); u.shadow.inherit=False
text(s,PAD,4.65,6.6,1.4,"공식몰의 실제 룩북 이미지를 기반으로 만든\n인스타그램(2건)·네이버 블로그·유튜브 예시 콘텐츠입니다.",15,INK2,spacing=1.4)
text(s,PAD,6.7,6,0.3,"SUMMER 2026 · 주니어 패션 브랜드몰",11.5,INK3,bold=True)

# 2 브랜드 분석
s=slide(); brand(s); kicker(s,"BRAND ANALYSIS")
text(s,PAD,1.05,11.8,1.0,"믿고 입히는 고퀄리티, 가장 트렌디한 주니어 브랜드",26,INK,True)
imgs=[("awk-036.jpg","데일리 스트릿"),("awk-032.jpg","서핑 래시가드"),("awk-044.jpg","액티브 스포츠"),("awk-060.jpg","여름 비치 룩")]
iw=(SW-2*PAD-3*0.2)/4
for i,(f,cap) in enumerate(imgs):
    l=PAD+i*(iw+0.2)
    pic_cover(s,f"{AWK}/{f}",l,1.95,iw,2.35)
    text(s,l,4.36,iw,0.3,cap,11,INK2,bold=True,align=PP_ALIGN.CENTER)
cw=(SW-2*PAD-0.25)/2
card(s,PAD,4.95,cw,1.9,"기본 정보",
     "· 업종: 주니어(초등~10대) 패션 공식 브랜드몰\n· 운영사: 제이씨물산 · 공식몰 awkmall.com\n· 브랜드: 에어워크주니어 · 메종피치 · 에버라스트\n· 가격대: 여름 아이템 5,900~22,900원 (합리적)",bsize=12)
card(s,PAD+cw+0.25,4.95,cw,1.9,"톤 & 무드",
     "밝고 활기찬 · 스포티 스트릿 헤리티지 · 여름 액티브.\n아이는 신나게, 부모는 믿고.\n컬러: 화이트 / 코발트 블루 / 코랄 오렌지 / 블랙",bsize=12)

# 3 콘텐츠 전략
s=slide(); brand(s); kicker(s,"CONTENT STRATEGY")
text(s,PAD,1.05,11.8,1.0,"채널마다 다른 역할, 하나의 브랜드 무드",26,INK,True)
strat=[("INSTAGRAM","발견 & 저장","룩북형 + 시즌 저장각 2건으로 도달과 저장을 동시에. 브랜드 무드 각인 + 지금 시즌(여름방학) 반응.",CORAL),
       ("NAVER BLOG","검색 유입 & 전환","'초등 여름옷·아동복 추천' 검색에 답하는 롱폼. 사이즈·후기 정보로 구매까지 연결.",BLUE),
       ("YOUTUBE","무드 각인 & 확산","실제 룩북컷을 시네마틱 영상으로. 활기찬 여름 무드를 30초 안에 각인.",CORAL)]
cw=(SW-2*PAD-2*0.25)/3
for i,(n,ti,bo,col) in enumerate(strat):
    l=PAD+i*(cw+0.25)
    c=card(s,l,2.2,cw,2.7,ti,bo,n=n,ncolor=col,tsize=19)

# 4 인스타 POST 1 (룩북)
def insta_slide(title_ko, sub, cards_glob, n, caption, tagcol, taglabel):
    s=slide(); brand(s); kicker(s,"INSTAGRAM")
    pill(s,SW-PAD-2.0,0.6,taglabel,tagcol)
    text(s,PAD,1.02,9,0.5,title_ko,24,INK,True)
    text(s,PAD,1.62,11.8,0.4,sub,13,INK2)
    # 카드 필름스트립
    gap=0.14
    total=SW-2*PAD
    cw=(total-(n-1)*gap)/n
    ch=cw*1350/1080
    for i in range(1,n+1):
        l=PAD+(i-1)*(cw+gap)
        p=cards_glob.format(i)
        s.shapes.add_picture(p,I(l),I(2.25),I(cw),I(ch))
    # 캡션 박스
    capy=2.25+ch+0.2
    cbox=rrect(s,PAD,capy,SW-2*PAD,SH-capy-0.35,fill=SOFT,border=LINE)
    tf=cbox.text_frame; tf.word_wrap=True; tf.margin_left=I(0.24);tf.margin_right=I(0.24);tf.margin_top=I(0.16);tf.margin_bottom=I(0.14)
    setp(tf.paragraphs[0],"발행 캡션",10.5,INK3,True,KFONT,PP_ALIGN.LEFT,1.0)
    p=tf.add_paragraph(); p.space_before=Pt(3); setp(p,caption,11.5,INK,False,KFONT,PP_ALIGN.LEFT,1.35)
    return s

cap1=("☀️ 여름은, 아이가 가장 활발한 계절 — 놀이터·물놀이터·주말 나들이, 매일 뛰어노는 아이에게 딱 맞는 여름 룩. "
      "데일리 스트릿 그래픽 티 / 냉감 액티브 트레이닝 / 서핑 래시가드 / 데님·버블 트렌디룩. 여름 아이템 5,900원부터. "
      "#에어워크주니어 #주니어패션 #아동복 #초등학생코디 #래시가드 #여름아동복")
insta_slide("POST 1 · 브랜딩 룩북 — 우리 아이의 여름 룩북","상시 브랜드 무드 각인용 에디토리얼 룩북 (6장)",
            f"{CON}/awk-carousel-{{}}.png",6,cap1,BLUE,"POST 1 · 룩북")

cap2=("🔖 저장 필수! 여름방학 필수템 5 — 방학 시작 전 이것만 챙기면 끝. "
      "01 물놀이 래시가드 · 02 냉감 트레이닝 · 03 데일리 팬츠 · 04 그래픽 티 · 05 여행 셋업룩. "
      "여름 아이템 5,900원부터, 지금 프로필 링크에서! #에어워크주니어 #여름방학 #방학룩 #방학필수템 #물놀이룩")
insta_slide("POST 2 · 시즌 저장각 — 여름방학 필수템 TOP 5","지금 시즌(여름방학)에 반응하는 카운트다운·체크리스트 (7장)",
            f"{CON}/awk2-vacation-{{}}.png",7,cap2,CORAL,"POST 2 · 저장각")

# 6 네이버 블로그
s=slide(); brand(s); kicker(s,"NAVER BLOG")
text(s,PAD,1.02,11.8,0.9,"검색으로 엄마를 부르는 롱폼",24,INK,True)
# 좌: 제목/본문 요약, 우: 이미지
lw=6.9
text(s,PAD,1.95,lw,0.3,"초등 여름옷 추천 · 아동복 브랜드",11,CORAL,bold=True)
text(s,PAD,2.28,lw,0.9,"초등학생 여름 옷 추천 | 에어워크주니어 AWK,\n물놀이부터 데일리까지 한 번에",18,INK,True,spacing=1.15)
body=("여름이 되면 아이 옷장이 제일 바쁘죠. 땀도 많고, 물놀이도 가야 하고, 요즘 아이들은 예쁘고 트렌디한 옷을 좋아하고요. "
      "그래서 요즘 자주 담는 곳이 에어워크주니어(AWK)입니다.\n\n"
      "· 고퀄리티인데 여름 아이템 5,900원부터 — 한 시즌 입고 크는 아이 옷에 부담이 적어요.\n"
      "· 데일리 스트릿: 스마일·바시티 그래픽 티 + 카고/데님 = 요즘 스타일.\n"
      "· 물놀이엔 서핑 래시가드 필수 — 자외선 차단 + 예쁜 배색.\n"
      "· 사이즈: 여유핏은 한 사이즈 업, 키·몸무게 기반 추천 제공.\n\n"
      "#에어워크주니어 #초등학생여름옷 #아동복추천 #주니어패션 #래시가드 #물놀이룩")
text(s,PAD,3.35,lw,3.6,body,12,INK2,spacing=1.4)
# 우측 이미지 2장
rx=PAD+lw+0.3; rw=SW-rx-PAD
pic_cover(s,f"{AWK}/awk-050.jpg",rx,1.95,rw,2.4)
pic_cover(s,f"{AWK}/awk-032.jpg",rx,4.5,rw,2.4)

# 7 유튜브
s=slide(); brand(s); kicker(s,"YOUTUBE")
text(s,PAD,1.02,11.8,0.9,"30초 안에 각인되는 여름 무드필름",24,INK,True)
text(s,PAD,1.62,11.8,0.4,"실제 룩북컷 4개를 Higgsfield로 시네마틱 영상화 (데일리→스트릿→물놀이→비치)",13,INK2)
sb=[("awk-011.jpg","DAILY"),("awk-036.jpg","STREET"),("awk-032.jpg","SWIM"),("awk-060.jpg","BEACH")]
iw=(SW-2*PAD-3*0.2)/4
for i,(f,tag) in enumerate(sb):
    l=PAD+i*(iw+0.2)
    pic_cover(s,f"{AWK}/{f}",l,2.2,iw,2.9)
    pill(s,l+0.12,2.32,tag,BLUE if i%2==0 else CORAL,size=9)
yb=rrect(s,PAD,5.35,SW-2*PAD,1.55,fill=SOFT,border=LINE)
tf=yb.text_frame; tf.word_wrap=True; tf.margin_left=I(0.24);tf.margin_top=I(0.16);tf.margin_right=I(0.24)
setp(tf.paragraphs[0],"🎬 영상 제목·설명·태그",10.5,INK3,True,KFONT,PP_ALIGN.LEFT,1.0)
p=tf.add_paragraph(); p.space_before=Pt(3); setp(p,"에어워크주니어 SUMMER 2026 | 우리 아이의 여름 룩북 (데일리·물놀이·스트릿)",13,INK,True,KFONT,PP_ALIGN.LEFT,1.2)
p=tf.add_paragraph(); p.space_before=Pt(3); setp(p,"가장 활발한 나이의 여름 ☀️ 데일리부터 물놀이 래시가드까지 · 5,900원부터 · awkmall.com   #에어워크주니어 #주니어패션 #아동복 #shorts",11,INK2,False,KFONT,PP_ALIGN.LEFT,1.3)
text(s,PAD,6.95,11.8,0.3,"※ 영상 파일(mp4)은 별도 전달 · 실제 집행 시 원본 고해상도로 교체 권장",10,INK3)

# 8 CLOSING
s=slide(bg=INK)
text(s,0,2.6,SW,0.4,"AWK · 에어워크주니어",14,CORAL,True,KFONT,PP_ALIGN.CENTER)
text(s,0,3.15,SW,1.4,"아이의 여름을, 가장 트렌디하게",34,WHITE,True,KFONT,PP_ALIGN.CENTER)
text(s,0,4.55,SW,0.8,"인스타그램 2건 · 네이버 블로그 · 유튜브 무드필름 — 실제 룩북 기반 예시 콘텐츠",14,RGBColor(0xC7,0xC7,0xD2),False,KFONT,PP_ALIGN.CENTER)
text(s,0,6.5,SW,0.3,"PSLAB · SNS 채널 자동화 솔루션",11,RGBColor(0x9A,0x9A,0xA6),True,KFONT,PP_ALIGN.CENTER)

OUT=f"{SP}/AWK-에어워크주니어-제안서.pptx"
prs.save(OUT)
print("saved", OUT, f"{os.path.getsize(OUT)//1024} KB slides:", len(prs.slides._sldIdLst))
