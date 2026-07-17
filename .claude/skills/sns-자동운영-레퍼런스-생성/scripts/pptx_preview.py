#!/usr/bin/env python3
# Render generated .pptx to an SVG grid (visual sanity check; LibreOffice is broken in this sandbox)
import base64
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

SP="/tmp/claude-0/-home-user-pslab/e949e4d9-ac40-56b3-bb03-0a51ce53f94b/scratchpad"
p=Presentation(f"{SP}/deck.pptx")
EMU=914400.0
SWpx=1280.0
SCALE=SWpx/(p.slide_width/EMU)   # px per inch
def px(v): return round(v/EMU*SCALE,1)
Wp=px(p.slide_width); Hp=px(p.slide_height)

def esc(t):
    return t.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")

def rgb_of(color):
    try:
        c=color.rgb
        return f"#{c}"
    except Exception:
        return None

def shape_fill(sh):
    try:
        f=sh.fill
        if f.type is not None and f.fore_color and f.fore_color.type is not None:
            c=f.fore_color.rgb
            return f"#{c}"
    except Exception:
        pass
    return None

def alpha_of(sh):
    # detect a:alpha inside solidFill
    try:
        el=sh.fill._xPr
        a=el.find(qn('a:solidFill'))
        if a is not None:
            clr=a.find(qn('a:srgbClr'))
            if clr is not None:
                al=clr.find(qn('a:alpha'))
                if al is not None:
                    return int(al.get('val'))/100000.0
    except Exception:
        pass
    return 1.0

ANCHOR={MSO_ANCHOR.TOP:'top',MSO_ANCHOR.MIDDLE:'middle',MSO_ANCHOR.BOTTOM:'bottom',None:'top'}
ALIGN={PP_ALIGN.LEFT:'start',PP_ALIGN.CENTER:'middle',PP_ALIGN.RIGHT:'end',None:'start'}

svgs=[]
for idx,s in enumerate(p.slides,1):
    parts=[f'<svg xmlns="http://www.w3.org/2000/svg" width="{Wp}" height="{Hp}" viewBox="0 0 {Wp} {Hp}" style="display:block">']
    for sh in s.shapes:
        if sh.left is None: continue
        l,t,w,h=px(sh.left),px(sh.top),px(sh.width),px(sh.height)
        # picture
        if sh.shape_type==13 or (hasattr(sh,'image') and getattr(sh,'image',None) is not None):
            try:
                blob=sh.image.blob; ct=sh.image.content_type
                b64=base64.b64encode(blob).decode()
                parts.append(f'<image x="{l}" y="{t}" width="{w}" height="{h}" preserveAspectRatio="xMidYMid slice" href="data:{ct};base64,{b64}"/>')
                continue
            except Exception:
                pass
        # filled shape
        fill=shape_fill(sh)
        op=alpha_of(sh)
        if fill:
            rx=8 if 'Rounded' in (sh.name or '') else 2
            parts.append(f'<rect x="{l}" y="{t}" width="{w}" height="{h}" rx="{rx}" fill="{fill}" fill-opacity="{op:.2f}"/>')
        # text
        if sh.has_text_frame and sh.text_frame.text.strip():
            tf=sh.text_frame
            anc=ANCHOR.get(tf.vertical_anchor,'top')
            # count lines to compute block height
            paras=[pp for pp in tf.paragraphs]
            # estimate line heights
            sizes=[]
            for pp in paras:
                sz=12
                for r in pp.runs:
                    if r.font.size: sz=r.font.size.pt
                sizes.append(sz)
            lineheights=[sz*1.25*SCALE/72 for sz in sizes]
            total=sum(lineheights)
            if anc=='middle': cy=t+h/2-total/2
            elif anc=='bottom': cy=t+h-total
            else: cy=t
            y=cy
            for pp,lh,sz in zip(paras,lineheights,sizes):
                al=ALIGN.get(pp.alignment,'start')
                if al=='middle': x=l+w/2
                elif al=='end': x=l+w
                else: x=l
                # concat runs
                txt="".join(r.text for r in pp.runs) or pp.text
                if not txt.strip():
                    y+=lh; continue
                col="#F4F6FB"; bold=False; fam="sans-serif"
                if pp.runs:
                    r0=pp.runs[0]
                    if r0.font.color and r0.font.color.type is not None:
                        col=rgb_of(r0.font.color) or col
                    bold=bool(r0.font.bold)
                    if r0.font.name and 'Georgia' in r0.font.name: fam="Georgia, serif"
                fpx=sz*SCALE/72
                fw="700" if bold else "400"
                baseline=y+fpx*0.82
                parts.append(f'<text x="{x}" y="{baseline:.1f}" font-family="{fam}" font-size="{fpx:.1f}" font-weight="{fw}" fill="{col}" text-anchor="{al}">{esc(txt)}</text>')
                y+=lh
    parts.append('</svg>')
    svgs.append((idx,"".join(parts)))

# assemble HTML grid
cells="".join(f'<figure><figcaption>slide {i}</figcaption><div class="fr">{svg}</div></figure>' for i,svg in svgs)
html=f'''<!doctype html><meta charset=utf8><style>
body{{margin:0;background:#333;font-family:sans-serif}}
.grid{{display:grid;grid-template-columns:1fr 1fr;gap:14px;padding:14px}}
figure{{margin:0}}
figcaption{{color:#ccc;font-size:12px;margin:0 0 4px}}
.fr{{outline:1px solid #000}}
svg{{width:100%;height:auto}}
</style><div class="grid">{cells}</div>'''
open(f"{SP}/ppt_preview.html","w").write(html)
print("wrote ppt_preview.html with",len(svgs),"slides")
